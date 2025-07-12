# universal_ingest.py - The Ultimate Document Devouring Machine
import asyncio
import os
from pathlib import Path
from typing import List, Dict, Any, Optional
from datetime import datetime
import mimetypes
import re

from ..core.extractors import SupabaseDocumentExtractor
from dotenv import load_dotenv

class UniversalDocumentProcessor:
    """
    Universal document processor - handles EVERYTHING
    Like having a Swiss Army knife for documents
    """
    
    def __init__(self):
        self.supported_extensions = {
            # Text-based files
            '.txt': 'text',
            '.md': 'markdown', 
            '.markdown': 'markdown',
            '.rst': 'text',
            '.rtf': 'rich_text',
            
            # Microsoft Office
            '.docx': 'word',
            '.doc': 'word_legacy',
            '.xlsx': 'excel',
            '.xls': 'excel_legacy', 
            '.pptx': 'powerpoint',
            '.ppt': 'powerpoint_legacy',
            
            # PDFs
            '.pdf': 'pdf',
            
            # Web formats
            '.html': 'html',
            '.htm': 'html',
            '.xml': 'xml',
            '.json': 'json',
            '.csv': 'csv',
            
            # Code files (great for technical documentation)
            '.py': 'python',
            '.js': 'javascript', 
            '.ts': 'typescript',
            '.java': 'java',
            '.cpp': 'cpp',
            '.c': 'c',
            '.sql': 'sql',
            '.yaml': 'yaml',
            '.yml': 'yaml',
            
            # Others
            '.epub': 'epub',
            '.mobi': 'mobi'
        }
        
        self._install_dependencies()
    
    def _install_dependencies(self):
        """Check and install required dependencies"""
        required_packages = [
            'python-docx',  # Word documents
            'openpyxl',     # Excel files
            'PyPDF2',       # PDF files
            'pdfplumber',   # Better PDF extraction
            'python-pptx',  # PowerPoint files
            'beautifulsoup4', # HTML parsing
            'ebooklib',     # EPUB files
            'pandas'        # CSV handling
        ]
        
        print("🔧 Checking document processing dependencies...")
        
        for package in required_packages:
            try:
                __import__(package.replace('-', '_').replace('python_', ''))
            except ImportError:
                print(f"⚠️  {package} not installed. Install with: pip install {package}")
    
    async def extract_text_from_file(self, file_path: Path) -> Dict[str, Any]:
        """
        Extract text from any supported file type
        Returns: {title, content, metadata, extraction_method}
        """
        
        extension = file_path.suffix.lower()
        
        if extension not in self.supported_extensions:
            print(f"⚠️  Unsupported file type: {extension}")
            return None
        
        try:
            extraction_method = self.supported_extensions[extension]
            
            if extraction_method == 'text':
                return await self._extract_text_file(file_path)
            elif extraction_method in ['markdown']:
                return await self._extract_markdown_file(file_path)
            elif extraction_method == 'word':
                return await self._extract_word_file(file_path)
            elif extraction_method == 'excel':
                return await self._extract_excel_file(file_path)
            elif extraction_method == 'powerpoint':
                return await self._extract_powerpoint_file(file_path)
            elif extraction_method == 'pdf':
                return await self._extract_pdf_file(file_path)
            elif extraction_method == 'html':
                return await self._extract_html_file(file_path)
            elif extraction_method == 'json':
                return await self._extract_json_file(file_path)
            elif extraction_method == 'csv':
                return await self._extract_csv_file(file_path)
            elif extraction_method in ['python', 'javascript', 'sql']:
                return await self._extract_code_file(file_path)
            elif extraction_method == 'epub':
                return await self._extract_epub_file(file_path)
            else:
                # Fallback to text extraction
                return await self._extract_text_file(file_path)
                
        except Exception as e:
            print(f"❌ Error processing {file_path}: {e}")
            return None
    
    async def _extract_text_file(self, file_path: Path) -> Dict[str, Any]:
        """Extract from plain text files"""
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            content = f.read()
        
        return {
            'title': file_path.stem.replace('_', ' ').replace('-', ' ').title(),
            'content': content,
            'metadata': {
                'extraction_method': 'text',
                'word_count': len(content.split()),
                'character_count': len(content)
            }
        }
    
    async def _extract_markdown_file(self, file_path: Path) -> Dict[str, Any]:
        """Extract from markdown files with enhanced parsing"""
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            content = f.read()
        
        # Extract title from first H1
        title_match = re.search(r'^#\s+(.+)$', content, re.MULTILINE)
        title = title_match.group(1).strip() if title_match else file_path.stem.replace('_', ' ').title()
        
        # Count headers and structure
        headers = re.findall(r'^#+\s+(.+)$', content, re.MULTILINE)
        
        return {
            'title': title,
            'content': content,
            'metadata': {
                'extraction_method': 'markdown',
                'header_count': len(headers),
                'section_titles': headers[:10],  # First 10 headers
                'word_count': len(content.split())
            }
        }
    
    async def _extract_word_file(self, file_path: Path) -> Dict[str, Any]:
        """Extract from Word documents (.docx)"""
        try:
            from docx import Document
            
            doc = Document(file_path)
            
            # Extract title (from first paragraph or document properties)
            title = file_path.stem.replace('_', ' ').title()
            if doc.core_properties.title:
                title = doc.core_properties.title
            elif doc.paragraphs and doc.paragraphs[0].text.strip():
                title = doc.paragraphs[0].text.strip()[:100]
            
            # Extract all text
            content_parts = []
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    content_parts.append(paragraph.text.strip())
            
            content = '\n\n'.join(content_parts)
            
            # Extract metadata
            props = doc.core_properties
            
            return {
                'title': title,
                'content': content,
                'metadata': {
                    'extraction_method': 'word_docx',
                    'author': props.author or 'Unknown',
                    'created': str(props.created) if props.created else None,
                    'modified': str(props.modified) if props.modified else None,
                    'word_count': len(content.split()),
                    'paragraph_count': len([p for p in doc.paragraphs if p.text.strip()])
                }
            }
            
        except ImportError:
            print("❌ python-docx not installed. Install with: pip install python-docx")
            return None
        except Exception as e:
            print(f"❌ Error reading Word file {file_path}: {e}")
            return None
    
    async def _extract_pdf_file(self, file_path: Path) -> Dict[str, Any]:
        """Extract from PDF files"""
        try:
            import pdfplumber
            
            content_parts = []
            metadata = {}
            
            with pdfplumber.open(file_path) as pdf:
                # Extract metadata
                if pdf.metadata:
                    metadata.update({
                        'author': pdf.metadata.get('Author', 'Unknown'),
                        'creator': pdf.metadata.get('Creator', 'Unknown'),
                        'title_from_pdf': pdf.metadata.get('Title', ''),
                        'subject': pdf.metadata.get('Subject', ''),
                        'page_count': len(pdf.pages)
                    })
                
                # Extract text from all pages
                for page_num, page in enumerate(pdf.pages):
                    page_text = page.extract_text()
                    if page_text and page_text.strip():
                        content_parts.append(f"[Page {page_num + 1}]\n{page_text.strip()}")
            
            content = '\n\n'.join(content_parts)
            
            # Determine title
            title = metadata.get('title_from_pdf', '').strip()
            if not title:
                title = file_path.stem.replace('_', ' ').title()
            
            metadata.update({
                'extraction_method': 'pdf_pdfplumber',
                'word_count': len(content.split()),
                'character_count': len(content)
            })
            
            return {
                'title': title,
                'content': content,
                'metadata': metadata
            }
            
        except ImportError:
            print("❌ pdfplumber not installed. Install with: pip install pdfplumber")
            # Fallback to PyPDF2
            try:
                import PyPDF2
                
                content_parts = []
                with open(file_path, 'rb') as file:
                    pdf_reader = PyPDF2.PdfReader(file)
                    
                    for page_num, page in enumerate(pdf_reader.pages):
                        page_text = page.extract_text()
                        if page_text.strip():
                            content_parts.append(f"[Page {page_num + 1}]\n{page_text.strip()}")
                
                content = '\n\n'.join(content_parts)
                
                return {
                    'title': file_path.stem.replace('_', ' ').title(),
                    'content': content,
                    'metadata': {
                        'extraction_method': 'pdf_pypdf2',
                        'page_count': len(pdf_reader.pages),
                        'word_count': len(content.split())
                    }
                }
                
            except ImportError:
                print("❌ PyPDF2 not installed. Install with: pip install PyPDF2")
                return None
        except Exception as e:
            print(f"❌ Error reading PDF {file_path}: {e}")
            return None
    
    async def _extract_excel_file(self, file_path: Path) -> Dict[str, Any]:
        """Extract from Excel files"""
        try:
            import pandas as pd
            
            # Read all sheets
            excel_file = pd.ExcelFile(file_path)
            content_parts = []
            
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(file_path, sheet_name=sheet_name)
                
                # Convert dataframe to readable text
                sheet_content = f"=== Sheet: {sheet_name} ===\n"
                sheet_content += f"Columns: {', '.join(df.columns)}\n"
                sheet_content += f"Rows: {len(df)}\n\n"
                
                # Add sample data (first 10 rows)
                sheet_content += df.head(10).to_string()
                content_parts.append(sheet_content)
            
            content = '\n\n'.join(content_parts)
            
            return {
                'title': file_path.stem.replace('_', ' ').title(),
                'content': content,
                'metadata': {
                    'extraction_method': 'excel',
                    'sheet_count': len(excel_file.sheet_names),
                    'sheet_names': excel_file.sheet_names,
                    'total_rows': sum(len(pd.read_excel(file_path, sheet_name=sheet)) for sheet in excel_file.sheet_names)
                }
            }
            
        except ImportError:
            print("❌ pandas or openpyxl not installed. Install with: pip install pandas openpyxl")
            return None
        except Exception as e:
            print(f"❌ Error reading Excel file {file_path}: {e}")
            return None
    
    async def _extract_powerpoint_file(self, file_path: Path) -> Dict[str, Any]:
        """Extract from PowerPoint files"""
        try:
            from pptx import Presentation
            
            prs = Presentation(file_path)
            content_parts = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_content = f"=== Slide {slide_num} ===\n"
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_content += shape.text.strip() + "\n"
                
                if slide_content.strip() != f"=== Slide {slide_num} ===":
                    content_parts.append(slide_content)
            
            content = '\n\n'.join(content_parts)
            
            return {
                'title': file_path.stem.replace('_', ' ').title(),
                'content': content,
                'metadata': {
                    'extraction_method': 'powerpoint',
                    'slide_count': len(prs.slides),
                    'word_count': len(content.split())
                }
            }
            
        except ImportError:
            print("❌ python-pptx not installed. Install with: pip install python-pptx")
            return None
        except Exception as e:
            print(f"❌ Error reading PowerPoint file {file_path}: {e}")
            return None
    
    async def _extract_html_file(self, file_path: Path) -> Dict[str, Any]:
        """Extract from HTML files"""
        try:
            from bs4 import BeautifulSoup
            
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
            
            soup = BeautifulSoup(content, 'html.parser')
            
            # Extract title
            title = soup.title.string if soup.title else file_path.stem.replace('_', ' ').title()
            
            # Extract text content
            text_content = soup.get_text(separator='\n', strip=True)
            
            return {
                'title': title,
                'content': text_content,
                'metadata': {
                    'extraction_method': 'html',
                    'word_count': len(text_content.split()),
                    'has_title': bool(soup.title)
                }
            }
            
        except ImportError:
            print("❌ beautifulsoup4 not installed. Install with: pip install beautifulsoup4")
            return None
    
    async def _extract_json_file(self, file_path: Path) -> Dict[str, Any]:
        """Extract from JSON files"""
        try:
            import json
            
            with open(file_path, 'r', encoding='utf-8') as f:
                data = json.load(f)
            
            # Convert JSON to readable text
            content = json.dumps(data, indent=2, ensure_ascii=False)
            
            return {
                'title': file_path.stem.replace('_', ' ').title(),
                'content': content,
                'metadata': {
                    'extraction_method': 'json',
                    'json_keys': list(data.keys()) if isinstance(data, dict) else [],
                    'data_type': type(data).__name__
                }
            }
        except Exception as e:
            print(f"❌ Error reading JSON file {file_path}: {e}")
            return None
    
    async def _extract_csv_file(self, file_path: Path) -> Dict[str, Any]:
        """Extract from CSV files"""
        try:
            import pandas as pd
            
            df = pd.read_csv(file_path)
            
            # Create readable content
            content = f"CSV Data Summary:\n"
            content += f"Columns: {', '.join(df.columns)}\n"
            content += f"Rows: {len(df)}\n\n"
            content += "Sample Data:\n"
            content += df.head(10).to_string()
            
            return {
                'title': file_path.stem.replace('_', ' ').title(),
                'content': content,
                'metadata': {
                    'extraction_method': 'csv',
                    'column_count': len(df.columns),
                    'row_count': len(df),
                    'columns': list(df.columns)
                }
            }
        except Exception as e:
            print(f"❌ Error reading CSV file {file_path}: {e}")
            return None
    
    async def _extract_code_file(self, file_path: Path) -> Dict[str, Any]:
        """Extract from code files"""
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            content = f.read()
        
        # Count functions, classes, etc.
        functions = len(re.findall(r'def\s+\w+', content))
        classes = len(re.findall(r'class\s+\w+', content))
        
        return {
            'title': f"Code: {file_path.name}",
            'content': content,
            'metadata': {
                'extraction_method': 'code',
                'language': file_path.suffix[1:],
                'function_count': functions,
                'class_count': classes,
                'line_count': len(content.splitlines())
            }
        }
    
    async def _extract_epub_file(self, file_path: Path) -> Dict[str, Any]:
        """Extract from EPUB files"""
        try:
            import ebooklib
            from ebooklib import epub
            from bs4 import BeautifulSoup
            
            book = epub.read_epub(str(file_path))
            
            title = book.get_metadata('DC', 'title')[0][0] if book.get_metadata('DC', 'title') else file_path.stem
            
            content_parts = []
            for item in book.get_items():
                if item.get_type() == ebooklib.ITEM_DOCUMENT:
                    soup = BeautifulSoup(item.get_content(), 'html.parser')
                    text = soup.get_text(separator='\n', strip=True)
                    if text.strip():
                        content_parts.append(text)
            
            content = '\n\n'.join(content_parts)
            
            return {
                'title': title,
                'content': content,
                'metadata': {
                    'extraction_method': 'epub',
                    'word_count': len(content.split()),
                    'chapter_count': len(content_parts)
                }
            }
            
        except ImportError:
            print("❌ ebooklib not installed. Install with: pip install ebooklib")
            return None

class UniversalDocumentIngestion:
    """Universal document ingestion system"""
    
    def __init__(self, supabase_url: str, supabase_key: str):
        self.extractor = SupabaseDocumentExtractor(supabase_url, supabase_key)
        self.processor = UniversalDocumentProcessor()
    
    async def ingest_everything(self, folder_path: str, document_type: str = "universal") -> List[str]:
        """Ingest ALL supported document types from a folder"""
        
        folder = Path(folder_path)
        documents = []
        
        print("🔍 UNIVERSAL DOCUMENT SCAN")
        print("=" * 40)
        
        # Get all files in folder
        all_files = list(folder.rglob('*'))
        supported_files = [f for f in all_files if f.suffix.lower() in self.processor.supported_extensions]
        
        print(f"📁 Total files found: {len(all_files)}")
        print(f"📊 Supported files: {len(supported_files)}")
        print(f"🎯 File types: {set(f.suffix.lower() for f in supported_files)}")
        print()
        
        # Process each supported file
        for file_path in supported_files:
            print(f"🔄 Processing: {file_path.name}")
            
            extracted = await self.processor.extract_text_from_file(file_path)
            
            if extracted:
                # Add common metadata
                extracted['metadata'].update({
                    'file_size': file_path.stat().st_size,
                    'file_modified': datetime.fromtimestamp(file_path.stat().st_mtime).isoformat(),
                    'file_extension': file_path.suffix,
                    'source_file': str(file_path)
                })
                
                document = {
                    'title': extracted['title'],
                    'content': extracted['content'],
                    'document_type': document_type,
                    'source_file': str(file_path),
                    'metadata': extracted['metadata']
                }
                
                documents.append(document)
                print(f"✅ Extracted: {extracted['title'][:50]}...")
            else:
                print(f"❌ Failed to extract: {file_path.name}")
        
        print(f"\n🚀 Ingesting {len(documents)} documents...")
        doc_ids = await self.extractor.ingest_documents(documents)
        
        print(f"✅ Successfully ingested {len(doc_ids)} documents!")
        return doc_ids

async def main():
    """Run universal document ingestion"""
    
    load_dotenv()
    
    print("🚀 UNIVERSAL DOCUMENT INGESTION SYSTEM")
    print("=" * 50)
    print("📋 Supported formats:")
    
    processor = UniversalDocumentProcessor()
    for ext, method in processor.supported_extensions.items():
        print(f"   {ext} → {method}")
    
    print("\n" + "=" * 50)
    
    # Initialize universal ingestion
    ingestion = UniversalDocumentIngestion(
        os.getenv('SUPABASE_URL'),
        os.getenv('SUPABASE_KEY')
    )
    
    # Ingest everything from documents folder
    doc_ids = await ingestion.ingest_everything("./documents", "strategic")
    
    # Get final analytics
    analytics = await ingestion.extractor.get_document_analytics()
    print(f"\n📊 FINAL ANALYTICS:")
    print(f"   • Total Documents: {analytics['total_documents']}")
    print(f"   • Document Types: {list(analytics['document_types'].keys())}")
    print(f"   • Database Health: {analytics['database_health']}")
    
    print("\n🎯 Your universal document intelligence system is OPERATIONAL!")

if __name__ == "__main__":
    asyncio.run(main())